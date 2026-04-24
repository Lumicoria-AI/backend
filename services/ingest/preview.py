"""Preview-artifact generation for office / binary formats.

At ingest time the pipeline stores a lightweight, renderable artifact next
to the original blob (same bucket, `{s3_key}.preview.html` or `.preview.json`)
so the frontend can render without parsing the source file itself.

    DOCX → HTML via mammoth, sanitized with bleach
    PPTX → HTML grid (one card per slide) via python-pptx
    XLSX → {sheets: [{name, headers, rows}]} JSON via openpyxl
    Others → no artifact, preview endpoint falls back to presigned original

All deps are optional.  When a library is missing the generator returns
None and the preview endpoint falls back to a download link.
"""

from __future__ import annotations

import json
from typing import Any, Dict, List, Optional, Tuple

import structlog

logger = structlog.get_logger(__name__)


# Mime-type dispatch.
_DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
_PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
_XLSX_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
}


def preview_kind(mime: str) -> Optional[str]:
    if mime in _DOCX_MIMES:
        return "docx"
    if mime in _PPTX_MIMES:
        return "pptx"
    if mime in _XLSX_MIMES:
        return "xlsx"
    return None


def preview_artifact_key(s3_key: str, mime: str) -> Optional[str]:
    kind = preview_kind(mime)
    if kind in {"docx", "pptx"}:
        return f"{s3_key}.preview.html"
    if kind == "xlsx":
        return f"{s3_key}.preview.json"
    return None


def preview_artifact_content_type(mime: str) -> str:
    kind = preview_kind(mime)
    if kind == "xlsx":
        return "application/json; charset=utf-8"
    return "text/html; charset=utf-8"


# ── DOCX → HTML ────────────────────────────────────────────────────────


def _render_docx(file_path: str) -> Optional[bytes]:
    try:
        import mammoth  # type: ignore
    except ImportError:
        return None
    try:
        import bleach  # type: ignore
    except ImportError:
        bleach = None  # type: ignore

    try:
        with open(file_path, "rb") as fh:
            result = mammoth.convert_to_html(fh)
        html = result.value or ""
    except Exception as e:
        logger.warning("docx_preview_failed", error=str(e))
        return None

    if bleach is not None:
        allowed_tags = [
            "p", "br", "h1", "h2", "h3", "h4", "h5", "h6",
            "strong", "b", "em", "i", "u", "s",
            "ul", "ol", "li", "a", "img",
            "table", "thead", "tbody", "tr", "th", "td",
            "blockquote", "code", "pre", "hr", "span", "div",
        ]
        allowed_attrs = {"a": ["href", "title"], "img": ["src", "alt"], "*": ["class"]}
        html = bleach.clean(html, tags=allowed_tags, attributes=allowed_attrs, strip=True)

    doc = (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        "<style>body{font-family:system-ui,sans-serif;max-width:820px;"
        "margin:40px auto;padding:0 16px;line-height:1.6;color:#222}"
        "table{border-collapse:collapse;width:100%;margin:12px 0}"
        "th,td{border:1px solid #ddd;padding:6px 10px;text-align:left}"
        "img{max-width:100%;height:auto}"
        "</style></head><body>"
        f"{html}"
        "</body></html>"
    )
    return doc.encode("utf-8")


# ── PPTX → HTML grid ───────────────────────────────────────────────────


def _render_pptx(file_path: str) -> Optional[bytes]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return None
    try:
        pres = Presentation(file_path)
    except Exception as e:
        logger.warning("pptx_preview_failed", error=str(e))
        return None

    def _escape(s: str) -> str:
        return (
            s.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )

    slides_html: List[str] = []
    for idx, slide in enumerate(pres.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text or "" for run in para.runs).strip()
                if text:
                    parts.append(_escape(text))
        body = "<br>".join(parts) or "<em>(no text)</em>"
        slides_html.append(
            f"<section class=\"slide\"><header>Slide {idx}</header><div>{body}</div></section>"
        )

    html = (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        "<style>body{font-family:system-ui,sans-serif;background:#f6f7f9;margin:0;padding:24px}"
        ".grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(320px,1fr));gap:16px}"
        ".slide{background:#fff;border:1px solid #e2e5ea;border-radius:10px;"
        "padding:16px;min-height:200px;box-shadow:0 1px 2px rgba(0,0,0,.04)}"
        ".slide header{font-size:12px;color:#6b7380;text-transform:uppercase;"
        "letter-spacing:.05em;margin-bottom:8px}"
        ".slide div{font-size:14px;line-height:1.5;color:#1f2430}"
        "</style></head><body>"
        "<div class=\"grid\">"
        + "".join(slides_html)
        + "</div></body></html>"
    )
    return html.encode("utf-8")


# ── XLSX → JSON ────────────────────────────────────────────────────────


def _render_xlsx(file_path: str) -> Optional[bytes]:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        return None
    try:
        wb = load_workbook(file_path, data_only=True, read_only=True)
    except Exception as e:
        logger.warning("xlsx_preview_failed", error=str(e))
        return None

    # Cap to keep artifact small — large sheets are better browsed via export.
    MAX_ROWS_PER_SHEET = 500
    MAX_COLS_PER_SHEET = 60

    sheets: List[Dict[str, Any]] = []
    for ws in wb.worksheets:
        rows: List[List[Any]] = []
        headers: Optional[List[Any]] = None
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            trimmed = list(row[:MAX_COLS_PER_SHEET])
            if i == 0:
                headers = [h if h is not None else f"col_{j+1}" for j, h in enumerate(trimmed)]
                continue
            if all(v is None or v == "" for v in trimmed):
                continue
            rows.append([None if v is None else str(v) for v in trimmed])
            if len(rows) >= MAX_ROWS_PER_SHEET:
                break
        sheets.append({
            "name": ws.title,
            "headers": headers or [],
            "rows": rows,
            "truncated": len(rows) >= MAX_ROWS_PER_SHEET,
        })

    try:
        return json.dumps({"sheets": sheets}, ensure_ascii=False).encode("utf-8")
    except Exception as e:
        logger.warning("xlsx_preview_serialize_failed", error=str(e))
        return None


# ── Public entry point ─────────────────────────────────────────────────


def render_preview(file_path: str, mime: str) -> Optional[Tuple[bytes, str]]:
    """Return `(artifact_bytes, content_type)` for a supported format,
    or None when no preview can be generated (missing dep, wrong mime,
    or parse failure).  Callers upload the bytes to
    `preview_artifact_key(s3_key, mime)`.
    """
    kind = preview_kind(mime)
    if kind == "docx":
        data = _render_docx(file_path)
    elif kind == "pptx":
        data = _render_pptx(file_path)
    elif kind == "xlsx":
        data = _render_xlsx(file_path)
    else:
        return None
    if data is None:
        return None
    return data, preview_artifact_content_type(mime)
